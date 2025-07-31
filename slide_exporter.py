"""Alternative slide export method using python-pptx and Pillow."""
import os
import logging
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
import io

logger = logging.getLogger(__name__)

def export_slides_as_images(pptx_path, output_dir, base_name):
    """
    Export PowerPoint slides as images using python-pptx.
    This is a fallback method when LibreOffice is not available or fails.
    """
    try:
        prs = Presentation(pptx_path)
        exported_files = []
        
        for idx, slide in enumerate(prs.slides):
            slide_num = idx + 1
            
            # Create a placeholder image with slide content summary
            img = Image.new('RGB', (1024, 768), color='white')
            draw = ImageDraw.Draw(img)
            
            # Add slide number
            title_text = f"Slide {slide_num}"
            
            # Try to extract text content from the slide
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text.strip())
            
            # Create the image with slide information
            try:
                # Try to use a nice font
                font_title = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 48)
                font_content = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 24)
            except:
                # Fallback to default font
                font_title = ImageFont.load_default()
                font_content = ImageFont.load_default()
            
            # Draw title
            draw.text((50, 50), title_text, fill='black', font=font_title)
            
            # Draw slide content preview
            y_offset = 150
            for i, text in enumerate(slide_text[:5]):  # Show first 5 text elements
                if text:
                    # Truncate long text
                    display_text = text[:80] + "..." if len(text) > 80 else text
                    draw.text((50, y_offset), f"• {display_text}", fill='darkgray', font=font_content)
                    y_offset += 40
            
            # Add a border
            draw.rectangle([5, 5, 1019, 763], outline='lightgray', width=2)
            
            # Save the image
            filename = f"{base_name}_slide_{slide_num}.png"
            filepath = os.path.join(output_dir, filename)
            img.save(filepath, 'PNG')
            exported_files.append(filename)
            
            logger.info(f"Created placeholder image for slide {slide_num}")
        
        return exported_files
        
    except Exception as e:
        logger.error(f"Failed to export slides using python-pptx: {e}")
        return []