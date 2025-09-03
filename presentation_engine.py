from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json
import uuid
import logging
from typing import Dict, List, Tuple, Optional
import xml.etree.ElementTree as ET
from layout_engine import LayoutEngine

# Set up logger
logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# PresentationEngine
# ---------------------------------------------------------------------------

class PresentationEngine:
    def __init__(self):
        self.namespace = {
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
        }
        self.layout_engine = LayoutEngine()
    
    def parse_presentation(self, filepath: str) -> Dict:
        """Parse a PPTX file into a structured JSON format."""
        prs = Presentation(filepath)
        structure = {
            'slides': [],
            'metadata': {
                'total_slides': len(prs.slides),
                'slide_width': prs.slide_width,
                'slide_height': prs.slide_height
            }
        }
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_data = {
                'id': f'slide_{slide_idx}',
                'index': slide_idx,
                'shapes': []
            }
            
            for shape_idx, shape in enumerate(slide.shapes):
                shape_data = {
                    'id': f'slide_{slide_idx}_shape_{shape_idx}',
                    'type': shape.shape_type,
                    'name': shape.name,
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height
                }
                
                # Extract text content
                if hasattr(shape, 'text'):
                    shape_data['text'] = shape.text
                    shape_data['text_frame'] = self._parse_text_frame(shape.text_frame) if hasattr(shape, 'text_frame') else None
                
                # Extract table data
                if shape.has_table:
                    shape_data['table'] = self._parse_table(shape.table)
                
                # Store XML for editing
                shape_data['xml'] = shape.element.xml
                
                slide_data['shapes'].append(shape_data)
            
            structure['slides'].append(slide_data)
        
        return structure
    
    def _parse_text_frame(self, text_frame) -> Dict:
        """Parse text frame details."""
        return {
            'paragraphs': [
                {
                    'text': p.text,
                    'level': p.level,
                    'alignment': str(p.alignment) if p.alignment else None
                }
                for p in text_frame.paragraphs
            ]
        }
    
    def _parse_table(self, table) -> Dict:
        """Parse table data."""
        return {
            'rows': table.rows,
            'columns': table.columns,
            'data': [
                [cell.text for cell in row.cells]
                for row in table.rows
            ]
        }
    
    def create_presentation_from_plan(self, plan: List[Dict], filepath: str) -> None:
        """Create a new presentation from the AI-generated plan."""
        prs = Presentation()
        
        for slide_plan in plan:
            slide_layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add title
            if 'title' in slide_plan:
                title_shape = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.5), Inches(9), Inches(1)
                )
                title_frame = title_shape.text_frame
                title_frame.text = slide_plan['title']
                title_frame.paragraphs[0].font.size = Pt(32)
                title_frame.paragraphs[0].font.bold = True
            
            # Add content based on type
            if slide_plan.get('content_type') == 'table' and 'data' in slide_plan:
                self._add_table(slide, slide_plan['data'])
            elif 'content' in slide_plan:
                self._add_bullet_points(slide, slide_plan['content'])
        
        prs.save(filepath)
    
    def _add_bullet_points(self, slide, content: List[str]) -> None:
        """Add bullet points to a slide."""
        content_shape = slide.shapes.add_textbox(
            Inches(0.5), Inches(2), Inches(9), Inches(5)
        )
        text_frame = content_shape.text_frame
        
        for idx, item in enumerate(content):
            if idx == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.level = 0
    
    def _add_table(self, slide, data: List[List[str]]) -> None:
        """Add a table to a slide."""
        rows = len(data)
        cols = len(data[0]) if data else 0
        
        if rows == 0 or cols == 0:
            return
        
        left = Inches(0.5)
        top = Inches(2)
        width = Inches(9)
        height = Inches(rows * 0.8)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Fill table
        for r, row_data in enumerate(data):
            for c, cell_data in enumerate(row_data):
                table.cell(r, c).text = str(cell_data)
                
                # Style header row
                if r == 0:
                    cell = table.cell(r, c)
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(14)
    
    def edit_shape(self, filepath: str, shape_id: str, command: str, 
                   context_mode: str, structure: Dict, llm) -> Tuple[bool, Dict]:
        """Edit a shape using LLM-generated modifications."""
        prs = Presentation(filepath)
        
        # Find the shape
        slide_idx, shape_idx = self._parse_shape_id(shape_id)
        if slide_idx >= len(prs.slides):
            return False, structure
        
        slide = prs.slides[slide_idx]
        if shape_idx >= len(slide.shapes):
            return False, structure
        
        shape = slide.shapes[shape_idx]
        
        # Prepare context for LLM
        context = self._prepare_edit_context(shape, command, context_mode, structure)
        
        # Get LLM response
        llm_response = llm.generate_response(context)
        
        # Apply the edit
        success = self._apply_edit(shape, llm_response)
        
        if success:
            # Save and re-parse
            prs.save(filepath)
            new_structure = self.parse_presentation(filepath)
            return True, new_structure
        
        return False, structure
    
    def _parse_shape_id(self, shape_id: str) -> Tuple[int, int]:
        """Parse shape ID to get slide and shape indices."""
        parts = shape_id.split('_')
        slide_idx = int(parts[1])
        shape_idx = int(parts[3])
        return slide_idx, shape_idx
    
    def _prepare_edit_context(self, shape, command: str, context_mode: str, structure: Dict) -> str:
        """Prepare context for LLM based on context mode."""
        prompt = f"""You are an expert at editing PowerPoint presentations using XML.
        
User command: {command}

Current shape XML:
{shape.element.xml}

"""
        
        if context_mode == 'global':
            # Add global context
            all_text = []
            for slide in structure['slides']:
                for s in slide['shapes']:
                    if 'text' in s and s['text']:
                        all_text.append(f"Slide {slide['index']}: {s['text']}")
            
            prompt += f"""
Global presentation context (all text):
{chr(10).join(all_text)}
"""
        
        prompt += """
Please provide the modified XML for this shape that implements the requested change.
Return ONLY the XML, no explanations or markdown formatting.
The XML should start with <p:sp or similar PowerPoint XML tags.
"""
        
        return prompt
    
    def _apply_edit(self, shape, xml_string: str) -> bool:
        """Apply XML edit to a shape."""
        try:
            # Clean the XML string - remove markdown formatting if present
            cleaned_xml = xml_string.strip()
            
            # Remove markdown code blocks
            if cleaned_xml.startswith('```xml'):
                cleaned_xml = cleaned_xml[6:]
            elif cleaned_xml.startswith('```'):
                cleaned_xml = cleaned_xml[3:]
                
            if cleaned_xml.endswith('```'):
                cleaned_xml = cleaned_xml[:-3]
                
            cleaned_xml = cleaned_xml.strip()
            
            # Parse the new XML
            new_element = ET.fromstring(cleaned_xml)
            
            # Replace the shape's element
            parent = shape.element.getparent()
            index = parent.index(shape.element)
            parent.remove(shape.element)
            parent.insert(index, new_element)
            
            return True
        except Exception as e:
            print(f"Error applying edit: {e}")
            print(f"XML was: {xml_string[:200]}...")
            return False
    
    def build_from_structured_text(self, filepath: str, structured_text: str, llm) -> bool:
        """Build out an entire presentation from structured text."""
        try:
            # Create prompt for LLM
            prompt = f"""You are an expert presentation architect. 
Your job is to convert structured text into a precise JSON plan for slides.

STRUCTURED TEXT (user input):
{structured_text}

===================
Allowed content_block types
===================
• bullets            – {{ "type": "bullets", "items": ["Point 1", "Point 2"] }}
• text               – {{ "type": "text", "text": "Paragraph content" }}
• content_box        – {{ "type": "content_box", "title": "Box Title", "items": ["Item 1"] }}
• table              – {{ "type": "table", "data": [["H","V"]] }}
• process            – {{ "type": "process", "steps": ["Step 1","Step 2"] }}
• kpi                – {{ "type": "kpi", "value": "95%", "label": "Satisfaction" }}
• comparison         – {{ "type": "comparison", "data": {{ "left":{{}}, "right":{{}} }} }}
• pyramid            – {{ "type": "pyramid", "levels": ["Top","Mid","Base"] }}
• column             – {{ "type": "column", "title": "Header", "items": ["A","B"] }}
• attribution        – {{ "type": "attribution", "text": "- Name, Org" }}
• highlight          – {{ "type": "highlight", "title": "Key", "text": "Important" }}

===================
Best-practice guidelines
===================
1. 3-6 bullets per slide, each ≤ 12 words.  
2. Avoid duplicate bullets or slides.  
3. Maintain consistent, professional tone and tense.  
4. Use content_box when a title plus list of items exists.  
5. Use kpi for standalone numeric stats.  
6. Use table when data naturally fits rows/columns.  
7. Use process for sequential steps; timeline for dated milestones.  
8. Choose the minimal block type needed – no mixtures that conflict.  

===================
Output requirements
===================
• Produce a VALID JSON array only – NO markdown fences, text, or comments.  
• Each element must include: "slide_number", "title", and "content_blocks".  
• Slide numbers must start at 1 and be sequential.  

Example skeleton (do NOT wrap in markdown):
[
  {{
    "slide_number": 1,
    "title": "Title text",
    "content_blocks": [
      {{
        "type": "bullets",
        "items": ["Example point"]
      }}
    ]
  }}
]
"""
            
            # Get LLM response
            response = llm.generate_response(prompt)
            
            # Clean and parse JSON
            cleaned_response = response.strip()
            if cleaned_response.startswith('```json'):
                cleaned_response = cleaned_response[7:]
            elif cleaned_response.startswith('```'):
                cleaned_response = cleaned_response[3:]
            if cleaned_response.endswith('```'):
                cleaned_response = cleaned_response[:-3]
            
            import json
            
            # Add debugging for JSON parsing issues
            logger.info(f"Attempting to parse JSON response (length: {len(cleaned_response)})")
            logger.info(f"JSON preview: {cleaned_response[:200]}...")
            
            try:
                slide_plan = json.loads(cleaned_response.strip())
            except json.JSONDecodeError as e:
                logger.error(f"JSON parsing failed: {e}")
                logger.error(f"Raw response: {cleaned_response}")
                logger.error(f"Error at position {e.pos}: '{cleaned_response[max(0, e.pos-20):e.pos+20]}'")
                raise Exception(f"Invalid JSON response from LLM: {str(e)}")
            
            # Handle both single slide object and array of slides
            if isinstance(slide_plan, dict):
                # Single slide object - wrap it in an array
                slide_plan = [slide_plan]
            elif isinstance(slide_plan, list):
                logger.info(f"Processing {len(slide_plan)} slides from LLM")
            else:
                raise Exception(f"Unexpected JSON structure type: {type(slide_plan)}")
            
            # Load the uploaded presentation to preserve themes, layouts, and formatting
            prs = Presentation(filepath)
            logger.info(f"Loaded uploaded presentation: {filepath}")
            logger.info(f"Original presentation has {len(prs.slides)} slides")
            
            # Clear existing slides using the simplest approach - create new presentation but copy theme
            # Save the slide masters and layouts from the original presentation
            original_slide_masters = prs.slide_masters
            original_slide_layouts = [layout for master in original_slide_masters for layout in master.slide_layouts]
            
            # For now, use a safer approach: just delete slides in reverse order
            slides_to_remove = list(prs.slides)
            for slide in reversed(slides_to_remove):
                slide_rId = slide.rId
                prs.part.drop_rel(slide_rId)
            
            # Clear the slide ID list
            prs.slides._sldIdLst.clear()
            
            # Build each slide from the plan
            for slide_idx, slide_data in enumerate(slide_plan):
                try:
                    if isinstance(slide_data, str):
                        logger.warning(f"Skipping invalid slide data (string): {slide_data}")
                        continue
                    
                    logger.info(f"Building slide {slide_idx + 1}: {slide_data.get('title', 'No title')}")
                    
                    # Add slide number to data
                    slide_data['slide_number'] = slide_idx + 1
                    
                    # Use LLM to select the best layout for this slide
                    layout_id = self.layout_engine.analyze_content_for_layout(slide_data, llm)
                    logger.info(f"Selected layout: {layout_id}")

                    # ------------------------------------------------------------------
                    # Capacity-aware validation & auto-shrink
                    # ------------------------------------------------------------------
                    try:
                        if self._needs_shrink(slide_data, layout_id):
                            logger.info("Content exceeds capacity – requesting summarized version")
                            slide_data = self._summarize_to_fit(slide_data, layout_id, llm)
                    except Exception as cap_err:
                        logger.warning(f"Capacity check failed, proceeding without shrink: {cap_err}")
                    
                    # Use a blank layout from the original presentation
                    try:
                        slide_layout = prs.slide_layouts[5]  # Blank layout (common)
                    except IndexError:
                        try:
                            slide_layout = prs.slide_layouts[6]  # Sometimes blank is at index 6
                        except IndexError:
                            slide_layout = prs.slide_layouts[0]  # Fallback to first available layout
                    
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # Apply the selected layout with content
                    self.layout_engine.apply_layout(slide, layout_id, slide_data)
                    
                except Exception as e:
                    logger.error(f"Error creating slide {slide_idx + 1}: {e}")
                    import traceback
                    traceback.print_exc()
                    continue
            
            # Save the presentation
            prs.save(filepath)
            logger.info(f"Saved presentation to {filepath}")
            
            # Add delay to ensure file is fully written before preview generation
            import time
            time.sleep(2)
            
            # Verify file was written correctly
            import os
            if os.path.exists(filepath):
                file_size = os.path.getsize(filepath)
                logger.info(f"Presentation file saved successfully: {file_size} bytes")
            else:
                logger.error(f"Presentation file not found after save: {filepath}")
            
            return True
            
        except Exception as e:
            print(f"Error building presentation: {e}")
            import traceback
            traceback.print_exc()
            return False

    # ------------------------------------------------------------------
    # Capacity helpers
    # ------------------------------------------------------------------
    def _needs_shrink(self, slide_data: Dict, layout_id: str) -> bool:
        """
        Quick heuristic check if slide_data is likely to overflow the chosen layout.
        Only simple bullet/item/row counts handled for now.
        """
        caps = self.layout_engine.capacity_map.get(layout_id, {})
        blocks = slide_data.get("content_blocks", [])

        # bullets capacity (single column)
        if "bullets" in caps:
            bullet_total = sum(len(b.get("items", [])) for b in blocks if b.get("type") == "bullets")
            if bullet_total > caps["bullets"]:
                return True
        # boxes
        if "boxes" in caps and "itemsPerBox" in caps:
            content_boxes = [b for b in blocks if b.get("type") == "content_box"]
            if len(content_boxes) > caps["boxes"]:
                return True
            for cb in content_boxes:
                if len(cb.get("items", [])) > caps["itemsPerBox"]:
                    return True
        # steps / process
        if "steps" in caps:
            steps_blocks = [b for b in blocks if b.get("type") in ("process", "bullets")]
            steps = 0
            for sb in steps_blocks:
                steps += len(sb.get("steps", sb.get("items", [])))
            if steps > caps["steps"]:
                return True
        # rows in table
        if "rows" in caps:
            for b in blocks:
                if b.get("type") == "table":
                    if len(b.get("data", [])) > caps["rows"]:
                        return True
        return False

    def _summarize_to_fit(self, slide_data: Dict, layout_id: str, llm) -> Dict:
        """
        Ask the LLM to shrink / summarize content so that it fits within the
        capacity heuristics of the specified layout. Returns adjusted slide data.
        """
        caps = self.layout_engine.capacity_map.get(layout_id, {})
        prompt = f"""You are an expert presentation editor.
The slide below has been assigned the layout '{layout_id}' with capacity limits:
{json.dumps(caps, indent=2)}.
Trim or summarise the content so it fits **within these limits** while keeping key information.
Return ONLY the JSON object of shape:
{{
  "title": "...",
  "content_blocks": [ ... ]
}}
Do NOT change field names.
Current slide JSON:
{json.dumps({k: slide_data[k] for k in ('title','content_blocks')}, indent=2)}
"""
        try:
            response = llm.generate_response(prompt).strip()
            # remove code fences if any
            if response.startswith("```"):
                response = response.split("```", 2)[1] if "```" in response[3:] else response.strip("```")
            adjusted = json.loads(response)
            if "title" in adjusted and "content_blocks" in adjusted:
                slide_data["title"] = adjusted["title"]
                slide_data["content_blocks"] = adjusted["content_blocks"]
                logger.info("Slide content successfully trimmed to fit capacity")
        except Exception as e:
            logger.warning(f"Failed to shrink content: {e}")
        return slide_data