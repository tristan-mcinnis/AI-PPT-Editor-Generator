from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import json
import logging
from typing import Dict, List, Tuple, Optional
import re

logger = logging.getLogger(__name__)

class LayoutEngine:
    def __init__(self, layout_file: str = "slide_layouts.json"):
        """Initialize the Layout Engine with layout definitions."""
        self.layouts = self._load_layouts(layout_file)
        self.layout_map = {layout['id']: layout for layout in self.layouts['layouts']}
        
    def _load_layouts(self, layout_file: str) -> Dict:
        """Load layout definitions from JSON file."""
        try:
            with open(layout_file, 'r') as f:
                return json.load(f)
        except Exception as e:
            logger.error(f"Failed to load layout file: {e}")
            return {"layouts": [], "layoutSelectionRules": {"defaultLayout": "content_with_title"}}
    
    def analyze_content_for_layout(self, slide_content: Dict, llm) -> str:
        """Use LLM to analyze content and select the most appropriate layout."""
        prompt = f"""You are an expert at analyzing presentation content and selecting the most appropriate slide layout.

Analyze the following slide content and select the BEST layout from the available options.

SLIDE CONTENT:
Title: {slide_content.get('title', 'No title')}
Content Blocks: {json.dumps(slide_content.get('content_blocks', []), indent=2)}

AVAILABLE LAYOUTS:
"""
        
        # Add layout descriptions to prompt
        for layout in self.layouts['layouts']:
            prompt += f"\n- {layout['id']}: {layout['description']}"
            prompt += f"\n  Content triggers: {', '.join(layout['contentTriggers'])}"
        
        prompt += """

ANALYSIS GUIDELINES:
1. Look for keywords and content patterns that match the content triggers
2. Consider the structure of the content (e.g., comparisons need two columns)
3. If the slide has a quote, always use the quote layout
4. If there's a large statistic or percentage, use big_number
5. For process steps or workflows, use process_flow
6. For timelines or roadmaps, use timeline
7. For data tables, use table
8. For opening slides, use title_slide
9. For section breaks, use section_divider
10. Default to content_with_title for general content

Return ONLY the layout ID (e.g., "two_content_boxes") with no explanation."""
        
        try:
            response = llm.generate_response(prompt)
            layout_id = response.strip().strip('"').strip("'")
            
            # Validate the layout ID
            if layout_id in self.layout_map:
                logger.info(f"LLM selected layout: {layout_id}")
                return layout_id
            else:
                logger.warning(f"Invalid layout ID from LLM: {layout_id}")
                return self._fallback_layout_selection(slide_content)
                
        except Exception as e:
            logger.error(f"LLM layout selection failed: {e}")
            return self._fallback_layout_selection(slide_content)
    
    def _fallback_layout_selection(self, slide_content: Dict) -> str:
        """Fallback rule-based layout selection when LLM fails."""
        title = slide_content.get('title', '').lower()
        content_blocks = slide_content.get('content_blocks', [])
        num_blocks = len(content_blocks)
        
        # Check if it's a title slide (first slide or has subtitle)
        if slide_content.get('slide_number', 0) == 1 or 'subtitle' in slide_content:
            return 'title_slide'
        
        # Check for quote (look for quote marks in title)
        if '"' in slide_content.get('title', '') or '"' in slide_content.get('title', ''):
            return 'quote'
        
        # Check for big number - more precise check (number is prominent in title)
        title_words = slide_content.get('title', '').split()
        for word in title_words[:3]:  # Check first 3 words for prominence
            if re.search(r'^\d+\.?\d*%$', word):
                return 'big_number'
        
        # Check for KPI/dashboard keywords
        if any(trigger in title for trigger in ['dashboard', 'metrics', 'kpi', 'performance metrics']):
            return 'kpi_dashboard'
        
        # Check for full image/overlay keywords
        if any(trigger in title for trigger in ['impactful', 'dramatic', 'visual statement']):
            return 'full_image_overlay'
        
        # Check content blocks for specific types
        for block in content_blocks:
            if block.get('type') == 'table':
                return 'table'
            elif block.get('type') == 'process':
                return 'process_flow'
            elif block.get('type') == 'kpi':
                return 'kpi_dashboard'
            elif block.get('type') == 'pyramid':
                return 'pyramid_hierarchy'
            elif block.get('type') == 'comparison':
                return 'comparison_table'
        
        # Check for specific number of content blocks
        if num_blocks == 2 and block.get('type') == 'content_box':
            return 'two_content_boxes'
        elif num_blocks == 3 and any(b.get('type') == 'column' for b in content_blocks):
            return 'three_columns'
        elif num_blocks == 4:
            return 'four_content_grid'
        elif num_blocks == 5:
            return 'five_columns'
        elif num_blocks == 6:
            return 'six_box_grid'
        
        # Check for split/vertical content
        if any(trigger in title for trigger in ['split', 'dual view', 'two perspectives', 'left vs right']):
            return 'vertical_split'
        
        # Check for hierarchy/pyramid keywords
        if any(trigger in title for trigger in ['hierarchy', 'pyramid', 'levels', 'foundation']):
            return 'pyramid_hierarchy'
        
        # Check for circular/ecosystem keywords
        if any(trigger in title for trigger in ['ecosystem', 'hub', 'central', 'core']):
            return 'circular_diagram'
        
        # Check for comparison keywords
        if any(trigger in title for trigger in ['comparison table', 'versus table', 'feature comparison']):
            return 'comparison_table'
        
        # Check for timeline keywords
        if any(trigger in title for trigger in ['timeline', 'roadmap', 'phases', 'journey']):
            return 'timeline'
        
        # Check for process keywords
        if any(trigger in title for trigger in ['process', 'steps', 'workflow', 'procedure']):
            return 'process_flow'
        
        # Check for alternating/dialogue keywords
        if any(trigger in title for trigger in ['alternating', 'back and forth', 'dialogue']):
            return 'alternating_content'
        
        # Check for before/after, problem/solution keywords
        if any(trigger in title for trigger in ['before after', 'then now', 'problem solution']):
            return 'top_bottom_split'
        
        # Check for sidebar keywords
        if any(trigger in title for trigger in ['sidebar', 'with notes', 'plus context']):
            return 'content_with_sidebar'
        
        # Check for highlights/takeaways keywords
        if any(trigger in title for trigger in ['key takeaways', 'highlights', 'callouts']):
            return 'content_with_highlights'
        
        # Check for gallery/showcase keywords
        if any(trigger in title for trigger in ['gallery', 'showcase', 'examples', 'portfolio']):
            if 'image' in title or 'photo' in title or 'visual' in title:
                return 'picture_grid_2x2' if num_blocks <= 4 else 'three_column_images'
        
        # Check for centered/focus keywords
        if any(trigger in title for trigger in ['key message', 'central idea', 'main point', 'focus']):
            return 'centered_content'
        
        # Check for vs/versus/comparison
        if any(trigger in title for trigger in ['vs', 'versus', 'comparison', 'highlights']):
            # Check if we have two content boxes
            if len([b for b in content_blocks if b.get('type') == 'content_box']) >= 2:
                return 'two_content_boxes'
            else:
                return 'two_column_text'
        
        # Check for next steps/CTA
        if any(trigger in title for trigger in ['next steps', 'action items', 'get started', 'conclusion']):
            return 'call_to_action'
        
        # Check for section divider
        if any(trigger in title for trigger in ['strategic', 'initiatives', 'chapter', 'section']) and not content_blocks:
            return 'section_divider'
        
        # Check for image/visual content
        if any(trigger in title for trigger in ['image', 'photo', 'visual', 'picture']):
            return 'left_content_right_image'
        
        # Default based on number of blocks
        if num_blocks == 2:
            return 'two_column_text'
        elif num_blocks >= 3 and num_blocks <= 5:
            return 'three_columns'
        
        # Default
        return self.layouts['layoutSelectionRules']['defaultLayout']
    
    def apply_layout(self, slide, layout_id: str, content: Dict) -> None:
        """Apply the selected layout to a slide with the given content."""
        layout = self.layout_map.get(layout_id)
        if not layout:
            logger.error(f"Layout {layout_id} not found")
            return
        
        logger.info(f"Applying layout: {layout_id}")
        
        # Apply each element in the layout
        for element in layout['elements']:
            self._create_element(slide, element, content)
    
    def _create_element(self, slide, element: Dict, content: Dict) -> None:
        """Create a specific element on the slide based on the layout definition."""
        elem_type = element['type']
        pos = element['position']
        style = element['style']
        
        if elem_type == 'title':
            self._add_title(slide, pos, style, content.get('title', ''))
            
        elif elem_type == 'subtitle':
            self._add_subtitle(slide, pos, style, content.get('subtitle', ''))
            
        elif elem_type == 'section_title':
            self._add_section_title(slide, pos, style, content.get('title', ''))
            
        elif elem_type == 'content':
            self._add_content(slide, pos, style, content.get('content_blocks', []))
            
        elif elem_type == 'content_box':
            # For two-column layouts, split content
            content_blocks = content.get('content_blocks', [])
            box_index = self._get_element_index(slide, elem_type)
            if box_index < len(content_blocks):
                self._add_content_box(slide, pos, style, content_blocks[box_index])
                
        elif elem_type == 'column':
            # For three-column layouts
            content_blocks = content.get('content_blocks', [])
            col_index = self._get_element_index(slide, elem_type)
            if col_index < len(content_blocks):
                self._add_column(slide, pos, style, content_blocks[col_index])
                
        elif elem_type == 'picture_placeholder':
            self._add_picture_placeholder(slide, pos, style)
            
        elif elem_type == 'quote_text':
            quote_content = self._extract_quote_content(content)
            self._add_quote(slide, pos, style, quote_content)
            
        elif elem_type == 'attribution':
            attribution = self._extract_attribution(content)
            self._add_attribution(slide, pos, style, attribution)
            
        elif elem_type == 'big_number':
            number = self._extract_big_number(content)
            self._add_big_number(slide, pos, style, number)
            
        elif elem_type == 'supporting_text':
            supporting = self._extract_supporting_text(content)
            self._add_supporting_text(slide, pos, style, supporting)
            
        elif elem_type == 'timeline':
            timeline_data = self._extract_timeline_data(content)
            self._add_timeline(slide, pos, style, timeline_data)
            
        elif elem_type == 'process_flow':
            process_data = self._extract_process_data(content)
            self._add_process_flow(slide, pos, style, process_data)
            
        elif elem_type == 'table':
            table_data = self._extract_table_data(content)
            self._add_table(slide, pos, style, table_data)
            
        elif elem_type == 'icon_list':
            list_items = self._extract_list_items(content)
            self._add_icon_list(slide, pos, style, list_items)
            
        elif elem_type == 'cta_title':
            self._add_cta_title(slide, pos, style, content.get('title', ''))
            
        elif elem_type == 'cta_content':
            cta_content = self._extract_cta_content(content)
            self._add_cta_content(slide, pos, style, cta_content)
            
        elif elem_type == 'flexible_content':
            self._add_flexible_content(slide, pos, style, content.get('content_blocks', []))
            
        # New element types for 20 additional layouts
        elif elem_type == 'text_column':
            content_blocks = content.get('content_blocks', [])
            col_index = self._get_element_index(slide, elem_type)
            if col_index < len(content_blocks):
                self._add_text_column(slide, pos, style, content_blocks[col_index])
                
        elif elem_type == 'grid_box':
            content_blocks = content.get('content_blocks', [])
            box_index = self._get_element_index(slide, elem_type)
            if box_index < len(content_blocks):
                self._add_grid_box(slide, pos, style, content_blocks[box_index])
                
        elif elem_type == 'main_content':
            self._add_main_content(slide, pos, style, content.get('content_blocks', []))
            
        elif elem_type == 'sidebar':
            # Get last content block for sidebar
            content_blocks = content.get('content_blocks', [])
            if content_blocks:
                self._add_sidebar(slide, pos, style, content_blocks[-1])
                
        elif elem_type == 'image_column':
            content_blocks = content.get('content_blocks', [])
            col_index = self._get_element_index(slide, elem_type)
            if col_index < len(content_blocks):
                self._add_image_column(slide, pos, style, content_blocks[col_index])
                
        elif elem_type == 'header_left':
            # Extract left header from split content
            headers = self._extract_split_headers(content)
            if headers and len(headers) > 0:
                self._add_header_left(slide, pos, style, headers[0])
                
        elif elem_type == 'header_right':
            headers = self._extract_split_headers(content)
            if headers and len(headers) > 1:
                self._add_header_right(slide, pos, style, headers[1])
                
        elif elem_type == 'narrow_column':
            content_blocks = content.get('content_blocks', [])
            col_index = self._get_element_index(slide, elem_type)
            if col_index < len(content_blocks):
                self._add_narrow_column(slide, pos, style, content_blocks[col_index])
                
        elif elem_type == 'image_with_caption':
            content_blocks = content.get('content_blocks', [])
            img_index = self._get_element_index(slide, elem_type)
            if img_index < len(content_blocks):
                self._add_image_with_caption(slide, pos, style, content_blocks[img_index])
                
        elif elem_type == 'large_image':
            self._add_large_image(slide, pos, style)
            
        elif elem_type == 'top_section':
            sections = self._extract_top_bottom_sections(content)
            if sections and 'top' in sections:
                self._add_top_section(slide, pos, style, sections['top'])
                
        elif elem_type == 'bottom_section':
            sections = self._extract_top_bottom_sections(content)
            if sections and 'bottom' in sections:
                self._add_bottom_section(slide, pos, style, sections['bottom'])
                
        elif elem_type == 'centered_box':
            # Use all content for centered box
            self._add_centered_box(slide, pos, style, {'text': self._extract_all_text(content)})
            
        elif elem_type == 'small_box':
            content_blocks = content.get('content_blocks', [])
            box_index = self._get_element_index(slide, elem_type)
            if box_index < len(content_blocks):
                self._add_small_box(slide, pos, style, content_blocks[box_index])
                
        elif elem_type == 'pyramid_top':
            pyramid_data = self._extract_pyramid_data(content)
            if pyramid_data and len(pyramid_data) > 0:
                self._add_pyramid_top(slide, pos, style, pyramid_data[0])
                
        elif elem_type == 'pyramid_middle':
            pyramid_data = self._extract_pyramid_data(content)
            if pyramid_data and len(pyramid_data) > 1:
                self._add_pyramid_middle(slide, pos, style, pyramid_data[1])
                
        elif elem_type == 'pyramid_bottom':
            pyramid_data = self._extract_pyramid_data(content)
            if pyramid_data and len(pyramid_data) > 2:
                self._add_pyramid_bottom(slide, pos, style, pyramid_data[2])
                
        elif elem_type == 'left_block':
            alternating_blocks = self._extract_alternating_blocks(content)
            block_index = self._get_element_index(slide, elem_type)
            if block_index < len(alternating_blocks['left']):
                self._add_left_block(slide, pos, style, alternating_blocks['left'][block_index])
                
        elif elem_type == 'right_block':
            alternating_blocks = self._extract_alternating_blocks(content)
            block_index = self._get_element_index(slide, elem_type)
            if block_index < len(alternating_blocks['right']):
                self._add_right_block(slide, pos, style, alternating_blocks['right'][block_index])
                
        elif elem_type == 'main_text':
            # Get main content blocks excluding highlights
            main_blocks = [b for b in content.get('content_blocks', []) if b.get('type') != 'highlight']
            self._add_main_text(slide, pos, style, main_blocks)
            
        elif elem_type == 'highlight_box':
            # Get highlight blocks
            highlight_blocks = [b for b in content.get('content_blocks', []) if b.get('type') == 'highlight']
            box_index = self._get_element_index(slide, elem_type)
            if box_index < len(highlight_blocks):
                self._add_highlight_box(slide, pos, style, highlight_blocks[box_index])
                
        elif elem_type == 'comparison_header':
            self._add_comparison_header(slide, pos, style, content.get('title', 'Comparison'))
            
        elif elem_type == 'comparison_left':
            comparison_data = self._extract_comparison_data(content)
            if comparison_data and 'left' in comparison_data:
                self._add_comparison_left(slide, pos, style, comparison_data['left'])
                
        elif elem_type == 'comparison_right':
            comparison_data = self._extract_comparison_data(content)
            if comparison_data and 'right' in comparison_data:
                self._add_comparison_right(slide, pos, style, comparison_data['right'])
                
        elif elem_type == 'center_circle':
            circular_data = self._extract_circular_data(content)
            if circular_data and 'center' in circular_data:
                self._add_center_circle(slide, pos, style, circular_data['center'])
                
        elif elem_type == 'satellite_box':
            circular_data = self._extract_circular_data(content)
            box_index = self._get_element_index(slide, elem_type)
            if circular_data and 'satellites' in circular_data and box_index < len(circular_data['satellites']):
                self._add_satellite_box(slide, pos, style, circular_data['satellites'][box_index])
                
        elif elem_type == 'full_image':
            self._add_full_image(slide, pos, style)
            
        elif elem_type == 'overlay_title':
            self._add_overlay_title(slide, pos, style, content.get('title', ''))
            
        elif elem_type == 'overlay_text':
            overlay_text = self._extract_overlay_text(content)
            self._add_overlay_text(slide, pos, style, overlay_text)
            
        elif elem_type == 'kpi_box':
            kpi_data = self._extract_kpi_data(content)
            box_index = self._get_element_index(slide, elem_type)
            if box_index < len(kpi_data):
                self._add_kpi_box(slide, pos, style, kpi_data[box_index])
                
        elif elem_type == 'chart_area':
            self._add_chart_area(slide, pos, style, {})
            
        elif elem_type == 'left_title':
            titles = self._extract_vertical_split_titles(content)
            if titles and 'left' in titles:
                self._add_left_title(slide, pos, style, titles['left'])
                
        elif elem_type == 'right_title':
            titles = self._extract_vertical_split_titles(content)
            if titles and 'right' in titles:
                self._add_right_title(slide, pos, style, titles['right'])
                
        elif elem_type == 'left_content':
            split_content = self._extract_vertical_split_content(content)
            if split_content and 'left' in split_content:
                self._add_left_content(slide, pos, style, split_content['left'])
                
        elif elem_type == 'right_content':
            split_content = self._extract_vertical_split_content(content)
            if split_content and 'right' in split_content:
                self._add_right_content(slide, pos, style, split_content['right'])
    
    def _get_element_index(self, slide, elem_type: str) -> int:
        """Get the index of the current element type being added."""
        # Simple counter - in real implementation, would track state
        return len([shape for shape in slide.shapes if hasattr(shape, '_element_type') and shape._element_type == elem_type])
    
    def _add_title(self, slide, pos: Dict, style: Dict, text: str) -> None:
        """Add a title to the slide."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        text_frame = shape.text_frame
        text_frame.text = text
        
        p = text_frame.paragraphs[0]
        p.font.size = Pt(style.get('fontSize', 32))
        if style.get('bold', False):
            p.font.bold = True
        if style.get('align') == 'center':
            p.alignment = PP_ALIGN.CENTER
    
    def _add_subtitle(self, slide, pos: Dict, style: Dict, text: str) -> None:
        """Add a subtitle to the slide."""
        if not text:
            text = ""  # Default subtitle if none provided
            
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        text_frame = shape.text_frame
        text_frame.text = text
        
        p = text_frame.paragraphs[0]
        p.font.size = Pt(style.get('fontSize', 24))
        if style.get('color'):
            p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        if style.get('align') == 'center':
            p.alignment = PP_ALIGN.CENTER
    
    def _add_section_title(self, slide, pos: Dict, style: Dict, text: str) -> None:
        """Add a section divider title."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        text_frame = shape.text_frame
        text_frame.text = text
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = text_frame.paragraphs[0]
        p.font.size = Pt(style.get('fontSize', 48))
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    
    def _add_content(self, slide, pos: Dict, style: Dict, content_blocks: List[Dict]) -> None:
        """Add content blocks to the slide."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        text_frame = shape.text_frame
        text_frame.margin_top = Inches(0.1)
        
        paragraph_idx = 0
        for block in content_blocks:
            if block['type'] == 'bullets' and 'items' in block:
                for idx, item in enumerate(block['items']):
                    if paragraph_idx == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = f"• {item}"
                    p.font.size = Pt(style.get('fontSize', 18))
                    p.level = 0
                    paragraph_idx += 1
                    
            elif block['type'] == 'text' and 'text' in block:
                if paragraph_idx == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = block['text']
                p.font.size = Pt(style.get('fontSize', 18))
                paragraph_idx += 1
    
    def _add_content_box(self, slide, pos: Dict, style: Dict, content_block: Dict) -> None:
        """Add a content box with title and items."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        # Add light background
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 240, 240)
        
        text_frame = shape.text_frame
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_right = Inches(0.2)
        text_frame.margin_top = Inches(0.2)
        text_frame.margin_bottom = Inches(0.2)
        
        # Add box title
        if content_block.get('title'):
            p = text_frame.paragraphs[0]
            p.text = content_block['title']
            p.font.size = Pt(18)
            p.font.bold = True
            
            # Add items
            if 'items' in content_block:
                for item in content_block['items']:
                    p = text_frame.add_paragraph()
                    p.text = f"• {item}"
                    p.font.size = Pt(14)
                    p.level = 1
    
    def _add_column(self, slide, pos: Dict, style: Dict, content_block: Dict) -> None:
        """Add a column of content."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        text_frame = shape.text_frame
        
        # Column title
        if content_block.get('title'):
            p = text_frame.paragraphs[0]
            p.text = content_block['title']
            p.font.size = Pt(16)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            
            # Column content
            if 'items' in content_block:
                for item in content_block['items']:
                    p = text_frame.add_paragraph()
                    p.text = item
                    p.font.size = Pt(12)
                    p.alignment = PP_ALIGN.LEFT
    
    def _add_picture_placeholder(self, slide, pos: Dict, style: Dict) -> None:
        """Add a picture placeholder."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        # Add border
        line = shape.line
        line.color.rgb = RGBColor(200, 200, 200)
        line.width = Pt(2)
        
        # Add placeholder text
        shape.text = "Image Placeholder"
        text_frame = shape.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(150, 150, 150)
    
    def _extract_quote_content(self, content: Dict) -> str:
        """Extract quote text from content."""
        # Look for quote in content blocks
        for block in content.get('content_blocks', []):
            if block.get('type') == 'quote':
                return block.get('text', '')
            # Also check for quotes in regular text
            if 'text' in block and ('"' in block['text'] or '"' in block['text']):
                return block['text']
        
        # Fallback to title if it contains quotes
        title = content.get('title', '')
        if '"' in title or '"' in title:
            return title
            
        return content.get('title', '')
    
    def _extract_attribution(self, content: Dict) -> str:
        """Extract attribution for quotes."""
        for block in content.get('content_blocks', []):
            if block.get('type') == 'attribution':
                return block.get('text', '')
            # Look for attribution patterns
            if 'text' in block:
                text = block['text']
                if text.startswith('-') or text.startswith('—') or 'said' in text.lower():
                    return text
        return ""
    
    def _extract_big_number(self, content: Dict) -> str:
        """Extract the big number/statistic from content."""
        import re
        
        # Look in title first
        title = content.get('title', '')
        numbers = re.findall(r'\d+\.?\d*%?', title)
        if numbers:
            return numbers[0]
        
        # Look in content blocks
        for block in content.get('content_blocks', []):
            if 'text' in block:
                numbers = re.findall(r'\d+\.?\d*%?', block['text'])
                if numbers:
                    return numbers[0]
        
        return "0"
    
    def _extract_supporting_text(self, content: Dict) -> str:
        """Extract supporting text for big numbers."""
        # Remove the number from title to get supporting text
        title = content.get('title', '')
        import re
        supporting = re.sub(r'\d+\.?\d*%?', '', title).strip()
        
        if not supporting:
            # Use first content block
            blocks = content.get('content_blocks', [])
            if blocks and 'text' in blocks[0]:
                supporting = blocks[0]['text']
        
        return supporting
    
    def _add_quote(self, slide, pos: Dict, style: Dict, quote_text: str) -> None:
        """Add a quote to the slide."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        text_frame = shape.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = text_frame.paragraphs[0]
        p.text = f'"{quote_text}"'
        p.font.size = Pt(style.get('fontSize', 28))
        p.font.italic = style.get('italic', True)
        p.alignment = PP_ALIGN.CENTER
    
    def _add_attribution(self, slide, pos: Dict, style: Dict, attribution: str) -> None:
        """Add attribution text."""
        if not attribution:
            return
            
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        text_frame = shape.text_frame
        
        p = text_frame.paragraphs[0]
        p.text = attribution
        p.font.size = Pt(style.get('fontSize', 18))
        p.alignment = PP_ALIGN.RIGHT
        p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    
    def _add_big_number(self, slide, pos: Dict, style: Dict, number: str) -> None:
        """Add a big number/statistic."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        text_frame = shape.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = text_frame.paragraphs[0]
        p.text = number
        p.font.size = Pt(style.get('fontSize', 72))
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(0x00, 0x66, 0xCC)
    
    def _add_supporting_text(self, slide, pos: Dict, style: Dict, text: str) -> None:
        """Add supporting text for big numbers."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        text_frame = shape.text_frame
        
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(style.get('fontSize', 24))
        p.alignment = PP_ALIGN.CENTER
    
    def _extract_timeline_data(self, content: Dict) -> List[Dict]:
        """Extract timeline data from content."""
        timeline_items = []
        
        for block in content.get('content_blocks', []):
            if block['type'] == 'bullets' and 'items' in block:
                for item in block['items']:
                    # Parse timeline items (e.g., "2024: Launch product")
                    if ':' in item:
                        parts = item.split(':', 1)
                        timeline_items.append({
                            'date': parts[0].strip(),
                            'event': parts[1].strip() if len(parts) > 1 else ''
                        })
                    else:
                        timeline_items.append({'date': '', 'event': item})
        
        return timeline_items
    
    def _add_timeline(self, slide, pos: Dict, style: Dict, timeline_data: List[Dict]) -> None:
        """Add a timeline to the slide."""
        if not timeline_data:
            return
        
        # Add timeline line
        line_y = pos['y'] + pos['height'] / 2
        line = slide.shapes.add_connector(
            1,  # Straight line
            Inches(pos['x']), Inches(line_y),
            Inches(pos['x'] + pos['width']), Inches(line_y)
        )
        line.line.width = Pt(3)
        line.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        
        # Add timeline items
        item_width = pos['width'] / len(timeline_data)
        for idx, item in enumerate(timeline_data):
            x_pos = pos['x'] + (idx * item_width) + (item_width / 2) - 0.5
            
            # Add milestone dot
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x_pos - 0.1), Inches(line_y - 0.1),
                Inches(0.2), Inches(0.2)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = RGBColor(0x00, 0x66, 0xCC)
            
            # Add date
            date_shape = slide.shapes.add_textbox(
                Inches(x_pos - 0.5), Inches(line_y - 0.8),
                Inches(1), Inches(0.5)
            )
            date_shape.text = item['date']
            date_shape.text_frame.paragraphs[0].font.size = Pt(12)
            date_shape.text_frame.paragraphs[0].font.bold = True
            date_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Add event
            event_shape = slide.shapes.add_textbox(
                Inches(x_pos - 0.75), Inches(line_y + 0.3),
                Inches(1.5), Inches(0.8)
            )
            event_shape.text = item['event']
            event_shape.text_frame.paragraphs[0].font.size = Pt(10)
            event_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            event_shape.text_frame.word_wrap = True
    
    def _extract_process_data(self, content: Dict) -> List[str]:
        """Extract process steps from content."""
        steps = []
        
        for block in content.get('content_blocks', []):
            if block['type'] == 'bullets' and 'items' in block:
                steps.extend(block['items'])
            elif block.get('type') == 'process' and 'steps' in block:
                steps.extend(block['steps'])
        
        return steps
    
    def _add_process_flow(self, slide, pos: Dict, style: Dict, steps: List[str]) -> None:
        """Add a process flow diagram."""
        if not steps:
            return
        
        num_steps = len(steps)
        step_width = (pos['width'] - 0.5 * (num_steps - 1)) / num_steps
        
        for idx, step in enumerate(steps):
            x_pos = pos['x'] + idx * (step_width + 0.5)
            
            # Add step box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x_pos), Inches(pos['y']),
                Inches(step_width), Inches(pos['height'])
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(230, 240, 250)
            
            # Add step text
            box.text = f"{idx + 1}. {step}"
            text_frame = box.text_frame
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            p = text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.CENTER
            p.word_wrap = True
            
            # Add arrow (except for last step)
            if idx < num_steps - 1:
                arrow_x = x_pos + step_width + 0.1
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW,
                    Inches(arrow_x), Inches(pos['y'] + pos['height']/2 - 0.15),
                    Inches(0.3), Inches(0.3)
                )
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = RGBColor(100, 100, 100)
    
    def _extract_table_data(self, content: Dict) -> List[List[str]]:
        """Extract table data from content."""
        for block in content.get('content_blocks', []):
            if block.get('type') == 'table' and 'data' in block:
                return block['data']
        
        # Try to create table from bullet points
        table_data = []
        for block in content.get('content_blocks', []):
            if block['type'] == 'bullets' and 'items' in block:
                for item in block['items']:
                    # Try to parse as table row (e.g., "Name | Value | Description")
                    if '|' in item:
                        row = [cell.strip() for cell in item.split('|')]
                        table_data.append(row)
        
        return table_data
    
    def _add_table(self, slide, pos: Dict, style: Dict, table_data: List[List[str]]) -> None:
        """Add a table to the slide."""
        if not table_data:
            return
        
        rows = len(table_data)
        cols = len(table_data[0]) if table_data else 0
        
        if rows == 0 or cols == 0:
            return
        
        table = slide.shapes.add_table(
            rows, cols,
            Inches(pos['x']), Inches(pos['y']),
            Inches(pos['width']), Inches(pos['height'])
        ).table
        
        # Style and fill table
        for r, row_data in enumerate(table_data):
            for c, cell_data in enumerate(row_data[:cols]):  # Ensure we don't exceed columns
                cell = table.cell(r, c)
                cell.text = str(cell_data)
                
                # Style header row
                if r == 0 and style.get('headerRow', True):
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0, 102, 204)
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(255, 255, 255)
                    paragraph.font.size = Pt(14)
                # Alternating row colors
                elif r > 0 and r % 2 == 0 and style.get('alternatingRows', True):
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    def _extract_list_items(self, content: Dict) -> List[Dict]:
        """Extract items for icon list."""
        items = []
        
        for block in content.get('content_blocks', []):
            if block['type'] == 'bullets' and 'items' in block:
                for item in block['items']:
                    items.append({'text': item, 'icon': '●'})  # Default icon
        
        return items
    
    def _add_icon_list(self, slide, pos: Dict, style: Dict, items: List[Dict]) -> None:
        """Add an icon list."""
        if not items:
            return
        
        item_height = pos['height'] / len(items)
        
        for idx, item in enumerate(items):
            y_pos = pos['y'] + idx * item_height
            
            # Add icon placeholder
            icon_size = style.get('iconSize', 0.5)
            icon_shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(pos['x']), Inches(y_pos),
                Inches(icon_size), Inches(icon_size)
            )
            icon_shape.fill.solid()
            icon_shape.fill.fore_color.rgb = RGBColor(0x00, 0x66, 0xCC)
            icon_shape.text = item.get('icon', '●')
            icon_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            icon_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Add text
            text_shape = slide.shapes.add_textbox(
                Inches(pos['x'] + icon_size + 0.2), Inches(y_pos),
                Inches(pos['width'] - icon_size - 0.2), Inches(item_height)
            )
            text_shape.text = item['text']
            text_shape.text_frame.paragraphs[0].font.size = Pt(style.get('fontSize', 16))
            text_shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    
    def _add_cta_title(self, slide, pos: Dict, style: Dict, text: str) -> None:
        """Add a call-to-action title."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        text_frame = shape.text_frame
        text_frame.text = text
        
        p = text_frame.paragraphs[0]
        p.font.size = Pt(style.get('fontSize', 40))
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    
    def _extract_cta_content(self, content: Dict) -> str:
        """Extract CTA content."""
        # Combine all content blocks for CTA
        cta_text = []
        for block in content.get('content_blocks', []):
            if block['type'] == 'text' and 'text' in block:
                cta_text.append(block['text'])
            elif block['type'] == 'bullets' and 'items' in block:
                cta_text.extend(block['items'])
        
        return '\n'.join(cta_text)
    
    def _add_cta_content(self, slide, pos: Dict, style: Dict, text: str) -> None:
        """Add CTA content."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        text_frame = shape.text_frame
        text_frame.text = text
        
        p = text_frame.paragraphs[0]
        p.font.size = Pt(style.get('fontSize', 20))
        p.alignment = PP_ALIGN.CENTER
    
    def _add_flexible_content(self, slide, pos: Dict, style: Dict, content_blocks: List[Dict]) -> None:
        """Add flexible content that adapts to content type."""
        # Similar to _add_content but with more flexibility
        self._add_content(slide, pos, style, content_blocks)
    
    # New layout element methods for the 20 additional layouts
    
    def _add_text_column(self, slide, pos: Dict, style: Dict, content_block: Dict) -> None:
        """Add a text column."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        text_frame = shape.text_frame
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        
        if content_block.get('title'):
            p = text_frame.paragraphs[0]
            p.text = content_block['title']
            p.font.size = Pt(18)
            p.font.bold = True
            
            if 'items' in content_block:
                for item in content_block['items']:
                    p = text_frame.add_paragraph()
                    p.text = f"• {item}"
                    p.font.size = Pt(style.get('fontSize', 16))
        elif 'text' in content_block:
            p = text_frame.paragraphs[0]
            p.text = content_block['text']
            p.font.size = Pt(style.get('fontSize', 16))
    
    def _add_grid_box(self, slide, pos: Dict, style: Dict, content_block: Dict) -> None:
        """Add a grid box with background."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        # Add background
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(245, 245, 245)
        
        # Add border
        line = shape.line
        line.color.rgb = RGBColor(200, 200, 200)
        line.width = Pt(1)
        
        text_frame = shape.text_frame
        text_frame.margin_all = Inches(0.15)
        
        if content_block.get('title'):
            p = text_frame.paragraphs[0]
            p.text = content_block['title']
            p.font.size = Pt(16)
            p.font.bold = True
            
            if 'items' in content_block:
                for item in content_block['items']:
                    p = text_frame.add_paragraph()
                    p.text = f"• {item}"
                    p.font.size = Pt(style.get('fontSize', 14))
    
    def _add_main_content(self, slide, pos: Dict, style: Dict, content_blocks: List[Dict]) -> None:
        """Add main content area."""
        self._add_content(slide, pos, style, content_blocks)
    
    def _add_sidebar(self, slide, pos: Dict, style: Dict, content_block: Dict) -> None:
        """Add a sidebar with background."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        # Add background color
        if style.get('backgroundColor'):
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(240, 240, 240)
        
        text_frame = shape.text_frame
        text_frame.margin_all = Inches(0.1)
        
        if 'text' in content_block:
            p = text_frame.paragraphs[0]
            p.text = content_block['text']
            p.font.size = Pt(style.get('fontSize', 14))
        elif 'items' in content_block:
            for idx, item in enumerate(content_block.get('items', [])):
                if idx == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(style.get('fontSize', 14))
    
    def _add_image_column(self, slide, pos: Dict, style: Dict, content_block: Dict) -> None:
        """Add an image column with caption."""
        # Image placeholder
        img_height = style.get('imageHeight', 3)
        img_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(img_height)
        )
        
        # Add border
        line = img_shape.line
        line.color.rgb = RGBColor(200, 200, 200)
        line.width = Pt(2)
        
        # Add placeholder text
        img_shape.text = "Image"
        img_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = img_shape.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(150, 150, 150)
        
        # Caption
        caption_y = pos['y'] + img_height + 0.1
        caption_height = style.get('captionHeight', 2) - 0.1
        
        caption_shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(caption_y),
            Inches(pos['width']), Inches(caption_height)
        )
        
        caption_text = content_block.get('caption', content_block.get('text', 'Caption'))
        caption_shape.text = caption_text
        caption_shape.text_frame.paragraphs[0].font.size = Pt(12)
        caption_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _add_header_left(self, slide, pos: Dict, style: Dict, text: str) -> None:
        """Add left header section."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        if style.get('backgroundColor'):
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(232, 244, 248)
        
        text_frame = shape.text_frame
        text_frame.text = text
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = text_frame.paragraphs[0]
        p.font.size = Pt(style.get('fontSize', 28))
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    
    def _add_header_right(self, slide, pos: Dict, style: Dict, text: str) -> None:
        """Add right header section."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        if style.get('backgroundColor'):
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(248, 232, 232)
        
        text_frame = shape.text_frame
        text_frame.text = text
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = text_frame.paragraphs[0]
        p.font.size = Pt(style.get('fontSize', 28))
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    
    def _add_narrow_column(self, slide, pos: Dict, style: Dict, content_block: Dict) -> None:
        """Add a narrow column for 5+ column layouts."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        text_frame = shape.text_frame
        
        if content_block.get('title'):
            p = text_frame.paragraphs[0]
            p.text = content_block['title']
            p.font.size = Pt(14)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            
            if 'items' in content_block:
                for item in content_block['items']:
                    p = text_frame.add_paragraph()
                    p.text = item
                    p.font.size = Pt(style.get('fontSize', 12))
    
    def _add_image_with_caption(self, slide, pos: Dict, style: Dict, content_block: Dict) -> None:
        """Add image with caption below."""
        # Calculate dimensions
        caption_height = 0.4
        image_height = pos['height'] - caption_height - 0.1
        
        # Image placeholder
        img_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(image_height)
        )
        
        # Style image placeholder
        line = img_shape.line
        line.color.rgb = RGBColor(200, 200, 200)
        line.width = Pt(1)
        
        img_shape.text = "Image"
        img_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = img_shape.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(150, 150, 150)
        
        # Caption
        caption_y = pos['y'] + image_height + 0.1
        caption_shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(caption_y),
            Inches(pos['width']), Inches(caption_height)
        )
        
        caption_text = content_block.get('caption', content_block.get('text', ''))
        caption_shape.text = caption_text
        caption_shape.text_frame.paragraphs[0].font.size = Pt(11)
        caption_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _add_large_image(self, slide, pos: Dict, style: Dict) -> None:
        """Add a large image placeholder."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        # Add border
        if style.get('border'):
            line = shape.line
            line.color.rgb = RGBColor(200, 200, 200)
            line.width = Pt(2)
        
        # Add placeholder text
        shape.text = "Large Image"
        text_frame = shape.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(150, 150, 150)
        p.font.size = Pt(18)
    
    def _add_top_section(self, slide, pos: Dict, style: Dict, content_block: Dict) -> None:
        """Add top section with background."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        if style.get('backgroundColor'):
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(245, 245, 245)
        
        text_frame = shape.text_frame
        text_frame.margin_all = Inches(0.2)
        
        if 'text' in content_block:
            p = text_frame.paragraphs[0]
            p.text = content_block['text']
            p.font.size = Pt(style.get('fontSize', 18))
        elif 'items' in content_block:
            for idx, item in enumerate(content_block.get('items', [])):
                if idx == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(style.get('fontSize', 18))
    
    def _add_bottom_section(self, slide, pos: Dict, style: Dict, content_block: Dict) -> None:
        """Add bottom section."""
        self._add_top_section(slide, pos, style, content_block)
    
    def _add_centered_box(self, slide, pos: Dict, style: Dict, content_block: Dict) -> None:
        """Add centered content box."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        if style.get('border'):
            line = shape.line
            line.color.rgb = RGBColor(200, 200, 200)
            line.width = Pt(2)
        
        text_frame = shape.text_frame
        text_frame.margin_all = Inches(0.3)
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        if 'text' in content_block:
            p = text_frame.paragraphs[0]
            p.text = content_block['text']
            p.font.size = Pt(style.get('fontSize', 20))
            p.alignment = PP_ALIGN.CENTER
    
    def _add_small_box(self, slide, pos: Dict, style: Dict, content_block: Dict) -> None:
        """Add small box for grid layouts."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        # Add light background and border
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(250, 250, 250)
        
        line = shape.line
        line.color.rgb = RGBColor(220, 220, 220)
        line.width = Pt(1)
        
        text_frame = shape.text_frame
        text_frame.margin_all = Inches(0.1)
        
        if content_block.get('title'):
            p = text_frame.paragraphs[0]
            p.text = content_block['title']
            p.font.size = Pt(14)
            p.font.bold = True
            
            if 'text' in content_block:
                p = text_frame.add_paragraph()
                p.text = content_block['text']
                p.font.size = Pt(style.get('fontSize', 12))
            elif 'items' in content_block:
                for item in content_block['items']:
                    p = text_frame.add_paragraph()
                    p.text = f"• {item}"
                    p.font.size = Pt(style.get('fontSize', 12))
    
    def _add_pyramid_top(self, slide, pos: Dict, style: Dict, text: str) -> None:
        """Add pyramid top level."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.TRAPEZOID,
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        # Style
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(68, 114, 196)
        
        # Text
        shape.text = text
        text_frame = shape.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(style.get('fontSize', 16))
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True
    
    def _add_pyramid_middle(self, slide, pos: Dict, style: Dict, text: str) -> None:
        """Add pyramid middle level."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.TRAPEZOID,
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(91, 155, 213)
        
        shape.text = text
        text_frame = shape.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(style.get('fontSize', 14))
        p.font.color.rgb = RGBColor(255, 255, 255)
    
    def _add_pyramid_bottom(self, slide, pos: Dict, style: Dict, text: str) -> None:
        """Add pyramid bottom level."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.TRAPEZOID,
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(112, 173, 71)
        
        shape.text = text
        text_frame = shape.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(style.get('fontSize', 14))
        p.font.color.rgb = RGBColor(255, 255, 255)
    
    def _add_left_block(self, slide, pos: Dict, style: Dict, content_block: Dict) -> None:
        """Add left-aligned block for alternating layout."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        text_frame = shape.text_frame
        if 'text' in content_block:
            p = text_frame.paragraphs[0]
            p.text = content_block['text']
            p.font.size = Pt(style.get('fontSize', 16))
            p.alignment = PP_ALIGN.LEFT
    
    def _add_right_block(self, slide, pos: Dict, style: Dict, content_block: Dict) -> None:
        """Add right-aligned block for alternating layout."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        text_frame = shape.text_frame
        if 'text' in content_block:
            p = text_frame.paragraphs[0]
            p.text = content_block['text']
            p.font.size = Pt(style.get('fontSize', 16))
            p.alignment = PP_ALIGN.RIGHT
    
    def _add_main_text(self, slide, pos: Dict, style: Dict, content_blocks: List[Dict]) -> None:
        """Add main text area."""
        self._add_content(slide, pos, style, content_blocks)
    
    def _add_highlight_box(self, slide, pos: Dict, style: Dict, content_block: Dict) -> None:
        """Add highlight box with colored background."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        # Background color based on style
        if style.get('backgroundColor'):
            fill = shape.fill
            fill.solid()
            color_map = {
                '#FFF2CC': RGBColor(255, 242, 204),  # Yellow
                '#E7F3FF': RGBColor(231, 243, 255),  # Blue
                '#FFE6E6': RGBColor(255, 230, 230)   # Red
            }
            rgb = color_map.get(style['backgroundColor'], RGBColor(240, 240, 240))
            fill.fore_color.rgb = rgb
        
        if style.get('border'):
            line = shape.line
            line.color.rgb = RGBColor(200, 200, 200)
            line.width = Pt(1)
        
        text_frame = shape.text_frame
        text_frame.margin_all = Inches(0.1)
        
        if content_block.get('title'):
            p = text_frame.paragraphs[0]
            p.text = content_block['title']
            p.font.size = Pt(12)
            p.font.bold = True
            
            if 'text' in content_block:
                p = text_frame.add_paragraph()
                p.text = content_block['text']
                p.font.size = Pt(style.get('fontSize', 14))
        elif 'text' in content_block:
            p = text_frame.paragraphs[0]
            p.text = content_block['text']
            p.font.size = Pt(style.get('fontSize', 14))
    
    def _add_comparison_header(self, slide, pos: Dict, style: Dict, text: str) -> None:
        """Add comparison table header."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        if style.get('backgroundColor'):
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(68, 114, 196)
        
        text_frame = shape.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(style.get('fontSize', 18))
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        if style.get('color') == 'white':
            p.font.color.rgb = RGBColor(255, 255, 255)
    
    def _add_comparison_left(self, slide, pos: Dict, style: Dict, content_block: Dict) -> None:
        """Add left comparison column."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        if style.get('backgroundColor'):
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(242, 242, 242)
        
        text_frame = shape.text_frame
        text_frame.margin_all = Inches(0.2)
        
        if content_block.get('title'):
            p = text_frame.paragraphs[0]
            p.text = content_block['title']
            p.font.size = Pt(18)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            
            if 'items' in content_block:
                for item in content_block['items']:
                    p = text_frame.add_paragraph()
                    p.text = f"• {item}"
                    p.font.size = Pt(style.get('fontSize', 16))
    
    def _add_comparison_right(self, slide, pos: Dict, style: Dict, content_block: Dict) -> None:
        """Add right comparison column."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        if style.get('backgroundColor'):
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(231, 230, 230)
        
        text_frame = shape.text_frame
        text_frame.margin_all = Inches(0.2)
        
        if content_block.get('title'):
            p = text_frame.paragraphs[0]
            p.text = content_block['title']
            p.font.size = Pt(18)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            
            if 'items' in content_block:
                for item in content_block['items']:
                    p = text_frame.add_paragraph()
                    p.text = f"• {item}"
                    p.font.size = Pt(style.get('fontSize', 16))
    
    def _add_center_circle(self, slide, pos: Dict, style: Dict, text: str) -> None:
        """Add center circle for circular diagram."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        if style.get('backgroundColor'):
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(68, 114, 196)
        
        shape.text = text
        text_frame = shape.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(style.get('fontSize', 18))
        p.font.bold = style.get('bold', True)
        
        if style.get('color') == 'white':
            p.font.color.rgb = RGBColor(255, 255, 255)
    
    def _add_satellite_box(self, slide, pos: Dict, style: Dict, content_block: Dict) -> None:
        """Add satellite box for circular diagram."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        # Add border
        line = shape.line
        line.color.rgb = RGBColor(200, 200, 200)
        line.width = Pt(1)
        
        text_frame = shape.text_frame
        text_frame.margin_all = Inches(0.1)
        
        if 'text' in content_block:
            p = text_frame.paragraphs[0]
            p.text = content_block['text']
            p.font.size = Pt(style.get('fontSize', 14))
            p.alignment = PP_ALIGN.CENTER
    
    def _add_full_image(self, slide, pos: Dict, style: Dict) -> None:
        """Add full slide image placeholder."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        # Semi-transparent
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(200, 200, 200)
        shape.fill.transparency = 0.7
        
        shape.text = "Full Background Image"
        text_frame = shape.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.font.size = Pt(24)
    
    def _add_overlay_title(self, slide, pos: Dict, style: Dict, text: str) -> None:
        """Add overlay title for full image slides."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(style.get('fontSize', 48))
        p.font.bold = style.get('bold', True)
        p.alignment = PP_ALIGN.CENTER
        
        # White text for overlay
        if style.get('color') == 'white':
            p.font.color.rgb = RGBColor(255, 255, 255)
    
    def _add_overlay_text(self, slide, pos: Dict, style: Dict, text: str) -> None:
        """Add overlay text for full image slides."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(style.get('fontSize', 24))
        p.alignment = PP_ALIGN.CENTER
        
        if style.get('color') == 'white':
            p.font.color.rgb = RGBColor(255, 255, 255)
    
    def _add_kpi_box(self, slide, pos: Dict, style: Dict, content_block: Dict) -> None:
        """Add KPI box with number and label."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        # Background color
        if style.get('backgroundColor'):
            fill = shape.fill
            fill.solid()
            color_map = {
                '#E7F3FF': RGBColor(231, 243, 255),
                '#FFE6E6': RGBColor(255, 230, 230),
                '#E6F4EA': RGBColor(230, 244, 234)
            }
            rgb = color_map.get(style['backgroundColor'], RGBColor(240, 240, 240))
            fill.fore_color.rgb = rgb
        
        text_frame = shape.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # KPI number
        p = text_frame.paragraphs[0]
        p.text = content_block.get('value', '0')
        p.font.size = Pt(style.get('numberSize', 36))
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # KPI label
        if content_block.get('label'):
            p = text_frame.add_paragraph()
            p.text = content_block['label']
            p.font.size = Pt(style.get('labelSize', 14))
            p.alignment = PP_ALIGN.CENTER
    
    def _add_chart_area(self, slide, pos: Dict, style: Dict, content_block: Dict) -> None:
        """Add chart area placeholder."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        if style.get('backgroundColor'):
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(245, 245, 245)
        
        shape.text = "Chart Area"
        text_frame = shape.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(150, 150, 150)
        p.font.size = Pt(18)
    
    def _add_left_title(self, slide, pos: Dict, style: Dict, text: str) -> None:
        """Add left title for vertical split."""
        self._add_title(slide, pos, style, text)
    
    def _add_right_title(self, slide, pos: Dict, style: Dict, text: str) -> None:
        """Add right title for vertical split."""
        self._add_title(slide, pos, style, text)
    
    def _add_left_content(self, slide, pos: Dict, style: Dict, content_blocks: List[Dict]) -> None:
        """Add left content for vertical split."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        if style.get('backgroundColor'):
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(248, 248, 248)
        
        # Add content
        text_frame = shape.text_frame
        text_frame.margin_all = Inches(0.2)
        
        for idx, block in enumerate(content_blocks):
            if block['type'] == 'bullets' and 'items' in block:
                for item_idx, item in enumerate(block['items']):
                    if idx == 0 and item_idx == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = f"• {item}"
                    p.font.size = Pt(style.get('fontSize', 16))
            elif block['type'] == 'text':
                if idx == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = block['text']
                p.font.size = Pt(style.get('fontSize', 16))
    
    def _add_right_content(self, slide, pos: Dict, style: Dict, content_blocks: List[Dict]) -> None:
        """Add right content for vertical split."""
        shape = slide.shapes.add_textbox(
            Inches(pos['x']), Inches(pos['y']), 
            Inches(pos['width']), Inches(pos['height'])
        )
        
        if style.get('backgroundColor'):
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(239, 239, 239)
        
        # Add content (similar to left_content)
        text_frame = shape.text_frame
        text_frame.margin_all = Inches(0.2)
        
        for idx, block in enumerate(content_blocks):
            if block['type'] == 'bullets' and 'items' in block:
                for item_idx, item in enumerate(block['items']):
                    if idx == 0 and item_idx == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = f"• {item}"
                    p.font.size = Pt(style.get('fontSize', 16))
            elif block['type'] == 'text':
                if idx == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = block['text']
                p.font.size = Pt(style.get('fontSize', 16))
    
    # Extraction methods for new layouts
    
    def _extract_split_headers(self, content: Dict) -> List[str]:
        """Extract headers for split header layout."""
        title = content.get("title", "")
        if " vs " in title:
            return title.split(" vs ", 1)
        elif " versus " in title:
            return title.split(" versus ", 1)
        elif " and " in title:
            return title.split(" and ", 1)
        elif ":" in title:
            return [part.strip() for part in title.split(":", 1)]
        else:
            # Default split
            words = title.split()
            mid = len(words) // 2
            return [" ".join(words[:mid]), " ".join(words[mid:])]
    
    def _extract_top_bottom_sections(self, content: Dict) -> Dict:
        """Extract top and bottom sections."""
        blocks = content.get("content_blocks", [])
        if len(blocks) >= 2:
            return {
                "top": blocks[0],
                "bottom": blocks[1] if len(blocks) > 1 else {"text": ""}
            }
        elif len(blocks) == 1:
            # Split single block
            return {
                "top": blocks[0],
                "bottom": {"text": ""}
            }
        return {"top": {"text": ""}, "bottom": {"text": ""}}
    
    def _extract_all_text(self, content: Dict) -> str:
        """Extract all text from content blocks."""
        texts = []
        for block in content.get("content_blocks", []):
            if block.get("type") == "text" and "text" in block:
                texts.append(block["text"])
            elif block.get("type") == "bullets" and "items" in block:
                texts.extend(block["items"])
        return "\n".join(texts)
    
    def _extract_pyramid_data(self, content: Dict) -> List[str]:
        """Extract pyramid hierarchy data."""
        pyramid_items = []
        for block in content.get("content_blocks", []):
            if block.get("type") == "pyramid":
                return block.get("levels", [])
            elif block.get("type") == "bullets" and "items" in block:
                # Use first 3 items for pyramid levels
                return block["items"][:3]
        
        # Fallback - use title words
        title = content.get("title", "")
        if title:
            return [title, "Middle Level", "Foundation"]
        return ["Top", "Middle", "Base"]
    
    def _extract_alternating_blocks(self, content: Dict) -> Dict:
        """Extract alternating left/right blocks."""
        blocks = content.get("content_blocks", [])
        left_blocks = []
        right_blocks = []
        
        for idx, block in enumerate(blocks):
            if idx % 2 == 0:
                left_blocks.append(block)
            else:
                right_blocks.append(block)
        
        return {"left": left_blocks, "right": right_blocks}
    
    def _extract_comparison_data(self, content: Dict) -> Dict:
        """Extract comparison data for comparison table."""
        blocks = content.get("content_blocks", [])
        
        # Look for explicit comparison blocks
        for block in blocks:
            if block.get("type") == "comparison":
                return block.get("data", {})
        
        # Otherwise split content in half
        if len(blocks) >= 2:
            return {
                "left": blocks[0],
                "right": blocks[1]
            }
        
        return {"left": {"title": "Option A", "items": []}, "right": {"title": "Option B", "items": []}}
    
    def _extract_circular_data(self, content: Dict) -> Dict:
        """Extract data for circular diagram."""
        center_text = content.get("title", "Core Concept")
        satellites = []
        
        for block in content.get("content_blocks", []):
            if block.get("type") == "bullets" and "items" in block:
                for item in block["items"][:4]:  # Max 4 satellites
                    satellites.append({"text": item})
            elif block.get("type") == "text":
                satellites.append({"text": block["text"]})
        
        # Ensure we have 4 satellites
        while len(satellites) < 4:
            satellites.append({"text": f"Element {len(satellites) + 1}"})
        
        return {
            "center": center_text,
            "satellites": satellites[:4]
        }
    
    def _extract_overlay_text(self, content: Dict) -> str:
        """Extract text for overlay."""
        blocks = content.get("content_blocks", [])
        if blocks and blocks[0].get("type") == "text":
            return blocks[0]["text"]
        return ""
    
    def _extract_kpi_data(self, content: Dict) -> List[Dict]:
        """Extract KPI data from content."""
        kpis = []
        
        # Look for KPI blocks
        for block in content.get("content_blocks", []):
            if block.get("type") == "kpi":
                kpis.append(block)
            elif block.get("type") == "bullets" and "items" in block:
                # Parse KPIs from bullets (e.g., "Revenue: $4.5M")
                for item in block["items"]:
                    if ":" in item:
                        label, value = item.split(":", 1)
                        kpis.append({"label": label.strip(), "value": value.strip()})
        
        # Ensure we have at least 3 KPIs
        while len(kpis) < 3:
            kpis.append({"label": f"Metric {len(kpis) + 1}", "value": "0"})
        
        return kpis
    
    def _extract_vertical_split_titles(self, content: Dict) -> Dict:
        """Extract titles for vertical split layout."""
        # Use split headers logic
        headers = self._extract_split_headers(content)
        if len(headers) >= 2:
            return {"left": headers[0], "right": headers[1]}
        return {"left": "Left Side", "right": "Right Side"}
    
    def _extract_vertical_split_content(self, content: Dict) -> Dict:
        """Extract content for vertical split layout."""
        blocks = content.get("content_blocks", [])
        
        # Split blocks evenly
        mid = len(blocks) // 2
        return {
            "left": blocks[:mid] if blocks else [],
            "right": blocks[mid:] if blocks else []
        }
